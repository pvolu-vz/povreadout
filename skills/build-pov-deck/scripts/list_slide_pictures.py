#!/usr/bin/env python3
"""List picture shapes on a single slide with their z-order index, EMU
position, EMU size, and the embedded image filename inside the package.

Used to author build_deck.py manifests — when a slide has multiple
pictures, the picture_index in the manifest is the z-order shown here.

EMU = English Metric Units. 914_400 EMU == 1 inch. Sizes are also shown
in inches for easier visual matching against PowerPoint's ruler.
"""
from __future__ import annotations

import argparse
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu


PICTURE_SHAPE_TYPE = 13
EMU_PER_INCH = 914_400


def main() -> int:
    skill_root = Path(__file__).resolve().parent.parent
    default_template = skill_root / "templates" / "pov-readout-template.pptx"

    ap = argparse.ArgumentParser(description=__doc__.split("\n\n")[0])
    ap.add_argument("slide_index", type=int, help="1-based slide number")
    ap.add_argument("--template", default=str(default_template))
    args = ap.parse_args()

    prs = Presentation(os.path.expanduser(args.template))
    if not (1 <= args.slide_index <= len(prs.slides)):
        sys.exit(
            f"slide_index={args.slide_index} out of range "
            f"(deck has {len(prs.slides)} slides)"
        )
    slide = prs.slides[args.slide_index - 1]
    pics = [s for s in slide.shapes if s.shape_type == PICTURE_SHAPE_TYPE]
    if not pics:
        print(f"slide={args.slide_index}: no pictures")
        return 0
    for i, pic in enumerate(pics):
        left_in = pic.left / EMU_PER_INCH
        top_in = pic.top / EMU_PER_INCH
        w_in = pic.width / EMU_PER_INCH
        h_in = pic.height / EMU_PER_INCH
        try:
            embedded = pic.image.filename or "<unnamed>"
        except Exception:
            embedded = "<unreadable>"
        print(
            f"slide={args.slide_index} picture_index={i} "
            f"pos=({left_in:.2f},{top_in:.2f})in "
            f"size=({w_in:.2f}x{h_in:.2f})in "
            f"embedded={embedded}"
        )
    return 0


if __name__ == "__main__":
    sys.exit(main())
