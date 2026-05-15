#!/usr/bin/env python3
"""List every slide in a template PPTX with its title and picture count.

Used to author build_deck.py manifests — you need to know the 1-based
slide_index for each section you intend to customize.

Output is one line per slide:
    slide=NN pictures=K title="..."

Use list_slide_pictures.py to drill into a specific slide's pictures.
"""
from __future__ import annotations

import argparse
import os
import sys
from pathlib import Path

from pptx import Presentation


PICTURE_SHAPE_TYPE = 13


def slide_title(slide) -> str:
    """First non-empty text on the slide. Falls back to '' if none."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            text = "".join(run.text for run in para.runs).strip()
            if text:
                return text
    return ""


def main() -> int:
    skill_root = Path(__file__).resolve().parent.parent
    default_template = skill_root / "templates" / "pov-readout-template.pptx"

    ap = argparse.ArgumentParser(description=__doc__.split("\n\n")[0])
    ap.add_argument("--template", default=str(default_template))
    args = ap.parse_args()

    prs = Presentation(os.path.expanduser(args.template))
    for i, slide in enumerate(prs.slides, 1):
        pic_count = sum(1 for s in slide.shapes if s.shape_type == PICTURE_SHAPE_TYPE)
        title = slide_title(slide)[:100]
        print(f'slide={i:02d} pictures={pic_count} title="{title}"')
    return 0


if __name__ == "__main__":
    sys.exit(main())
