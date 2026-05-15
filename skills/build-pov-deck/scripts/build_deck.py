#!/usr/bin/env python3
"""Clone a POV read-out template PPTX and apply a YAML manifest of edits.

Two kinds of edits are supported:
  - text_replacements: find/replace applied across every text frame
  - slide_overrides:  per-slide image swaps + per-slide text replacements

See ../SKILL.md for the manifest schema and intended workflow.

Why "remove + add_picture" for image swaps:
  python-pptx exposes Picture.image as read-only. Rewriting the underlying
  blip would require carefully manipulating the package's relationships —
  doable but easy to corrupt. Removing the shape and re-adding a fresh
  Picture at the same (left, top, width, height) is the simplest path
  that preserves layout fidelity and works on every PPTX we've seen.
  The trade-off: any shape-level effects (animations, alt text) on the
  swapped picture are lost. POV decks don't use those, so it's a clean win.
"""
from __future__ import annotations

import argparse
import os
import sys
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.util import Emu


PICTURE_SHAPE_TYPE = 13  # MSO_SHAPE_TYPE.PICTURE


def expand(p: str | Path) -> Path:
    return Path(os.path.expanduser(str(p))).resolve()


def slide_pictures(slide):
    """Return picture shapes on a slide in z-order (document order)."""
    return [s for s in slide.shapes if s.shape_type == PICTURE_SHAPE_TYPE]


def apply_text_replacements(slide, replacements) -> int:
    """Run-level first (preserves formatting), then paragraph-level fallback
    for text that got split across multiple runs by PowerPoint's autoformatter.
    Returns the number of substitutions made on this slide."""
    n = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        for para in tf.paragraphs:
            # Run-level pass
            for run in para.runs:
                for r in replacements:
                    if r["find"] in run.text:
                        run.text = run.text.replace(r["find"], r["replace"])
                        n += 1
            # Paragraph-level fallback for finds that span run boundaries.
            # When this fires, formatting in non-first runs is lost — that's
            # acceptable for our use case (customer-name swaps in titles).
            for r in replacements:
                joined = "".join(run.text for run in para.runs)
                if r["find"] in joined:
                    new_joined = joined.replace(r["find"], r["replace"])
                    if new_joined != joined and para.runs:
                        para.runs[0].text = new_joined
                        for extra in list(para.runs)[1:]:
                            extra.text = ""
                        n += 1
    return n


def swap_image(slide, picture_index: int, source: Path) -> tuple[int, int, int, int]:
    """Remove picture at z-order picture_index, replace with `source` at the
    same position+size. Returns (left, top, width, height) in EMU for logging."""
    pics = slide_pictures(slide)
    if picture_index >= len(pics):
        raise IndexError(
            f"picture_index={picture_index} out of range "
            f"(slide has {len(pics)} pictures)"
        )
    old = pics[picture_index]
    left, top, width, height = old.left, old.top, old.width, old.height
    sp = old._element
    sp.getparent().remove(sp)
    slide.shapes.add_picture(
        str(source), left, top, width=Emu(width), height=Emu(height)
    )
    return left, top, width, height


def remove_slide(prs, slide_index_0based: int) -> None:
    """Drop a slide from the deck. python-pptx has no public API for this,
    so we manipulate the sldIdLst and drop the part relationship.

    We do not delete the underlying slide part — leaving it orphaned is
    harmless (PowerPoint ignores unreferenced parts) and avoids the brittle
    work of garbage-collecting media/_rels chains.
    """
    sldIdLst = prs.slides._sldIdLst
    sld_ids = list(sldIdLst)
    if slide_index_0based >= len(sld_ids):
        raise IndexError(f"slide_index={slide_index_0based + 1} out of range")
    target = sld_ids[slide_index_0based]
    rId = target.get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )
    sldIdLst.remove(target)
    if rId:
        prs.part.drop_rel(rId)


def validate(manifest, prs, template_dir: Path) -> None:
    total = len(prs.slides)
    for ov in manifest.get("slide_overrides", []) or []:
        idx = ov.get("slide_index")
        if not isinstance(idx, int) or not (1 <= idx <= total):
            sys.exit(
                f"manifest error: slide_overrides entry has invalid "
                f"slide_index={idx!r} (deck has {total} slides)"
            )
        slide = prs.slides[idx - 1]
        pics = slide_pictures(slide)
        for img in ov.get("image_replacements", []) or []:
            pi = img.get("picture_index")
            if not isinstance(pi, int) or not (0 <= pi < len(pics)):
                sys.exit(
                    f"manifest error: slide {idx} has {len(pics)} picture(s); "
                    f"picture_index={pi!r} is out of range"
                )
            src = img.get("source")
            if not src:
                sys.exit(f"manifest error: slide {idx} image_replacement missing 'source'")
            srcp = expand(src)
            if not srcp.exists():
                sys.exit(f"manifest error: slide {idx} source not found: {srcp}")
    for ri in manifest.get("remove_slides", []) or []:
        if not isinstance(ri, int) or not (1 <= ri <= total):
            sys.exit(
                f"manifest error: remove_slides entry {ri!r} out of range "
                f"(deck has {total} slides)"
            )


def build(template: Path, manifest_path: Path, output: Path) -> int:
    if not template.exists():
        sys.exit(f"template not found: {template}")
    if not manifest_path.exists():
        sys.exit(f"manifest not found: {manifest_path}")

    manifest = yaml.safe_load(manifest_path.read_text()) or {}
    prs = Presentation(str(template))

    validate(manifest, prs, template.parent)

    global_text = manifest.get("text_replacements", []) or []
    if global_text:
        total_subs = 0
        for slide in prs.slides:
            total_subs += apply_text_replacements(slide, global_text)
        print(f"slide=* op=text count={total_subs} replacements={len(global_text)}")

    for ov in manifest.get("slide_overrides", []) or []:
        idx = ov["slide_index"]
        slide = prs.slides[idx - 1]
        for img in ov.get("image_replacements", []) or []:
            src = expand(img["source"])
            l, t, w, h = swap_image(slide, img["picture_index"], src)
            print(
                f"slide={idx} op=image picture_index={img['picture_index']} "
                f"src={src} size_emu=({w}x{h})"
            )
        local_text = ov.get("text_replacements", []) or []
        if local_text:
            n = apply_text_replacements(slide, local_text)
            print(f"slide={idx} op=text count={n} replacements={len(local_text)}")

    # Process removals last and in reverse order so earlier indices stay valid
    # as we mutate the slide list.
    for ri in sorted(manifest.get("remove_slides", []) or [], reverse=True):
        remove_slide(prs, ri - 1)
        print(f"slide={ri} op=remove")

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    print(f"wrote={output}")
    return 0


def main() -> int:
    skill_root = Path(__file__).resolve().parent.parent
    default_template = skill_root / "templates" / "pov-readout-template.pptx"

    ap = argparse.ArgumentParser(description=__doc__.split("\n\n")[0])
    ap.add_argument("--template", default=str(default_template))
    ap.add_argument("--manifest", required=True)
    ap.add_argument("--output", required=True)
    args = ap.parse_args()

    return build(
        expand(args.template),
        expand(args.manifest),
        expand(args.output),
    )


if __name__ == "__main__":
    sys.exit(main())
